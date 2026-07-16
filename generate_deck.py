"""
Archimedes AI — SAP LeanIX branded PowerPoint deck generator
"""

import os
import tempfile
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SKILL_DIR = "/Users/I519409/.claude/plugins/cache/btm-services-marketplace/leanix-docs/1.17.0/skills/generate-deck"
BASE_PATH = f"{SKILL_DIR}/templates/SAP_LeanIX_Base.pptx"
OUTPUT_PATH = "/Users/I519409/dev/archimedes-ai/Archimedes_AI_Presentation.pptx"

# SAP Brand Palette
DARK_BLUE = RGBColor(0x00, 0x2A, 0x86)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x00, 0x00, 0x00)
GRAY      = RGBColor(0x5B, 0x73, 0x8B)
TABLE_ALT = RGBColor(0xEA, 0xEC, 0xEE)
GREEN     = RGBColor(0x18, 0x7A, 0x38)

# Title layout constants
TITLE_LEFT   = 258554
TITLE_TOP    = 264435
TITLE_WIDTH  = 9277232
TITLE_HEIGHT = 369332
TITLE_HEIGHT_DIVIDER = 677108

# Logo constants
LOGO_LEFT   = 9981232
LOGO_TOP    = 269101
LOGO_WIDTH  = 1709244
LOGO_HEIGHT = 360000

# Table body area
PH_LEFT     = 504000
PH_BODY_TOP = 1548000
PH_WIDTH    = 11186477


def extract_logo(prs):
    """Extract blue SAP LeanIX logo PNG from donor slide 1."""
    slide = prs.slides[1]
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            img_bytes = shape.image.blob
            if len(img_bytes) < 15000:
                tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                tmp.write(img_bytes)
                tmp.close()
                return tmp.name
    return None


def add_logo(slide, logo_path):
    if logo_path:
        slide.shapes.add_picture(logo_path, LOGO_LEFT, LOGO_TOP, LOGO_WIDTH, LOGO_HEIGHT)


def set_title(slide, text, is_divider=False):
    ph = slide.placeholders[0]
    ph.left   = TITLE_LEFT
    ph.top    = TITLE_TOP
    ph.width  = TITLE_WIDTH
    ph.height = TITLE_HEIGHT_DIVIDER if is_divider else TITLE_HEIGHT
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE


def add_content_slide(prs, layout, logo_path, title, bullets, font_size=14):
    slide = prs.slides.add_slide(layout)
    set_title(slide, title)
    add_logo(slide, logo_path)
    ph = slide.placeholders[10]
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = "\u2022  " + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = BLACK
        p.space_before = Pt(5)
        p.space_after = Pt(5)
    return slide


def add_divider_slide(prs, layout, logo_path, title):
    slide = prs.slides.add_slide(layout)
    set_title(slide, title, is_divider=True)
    add_logo(slide, logo_path)
    return slide


def add_table_slide(prs, layout, logo_path, title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(layout)
    set_title(slide, title)
    add_logo(slide, logo_path)

    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)
    row_height = Emu(500000)
    total_height = row_height * num_rows

    tbl = slide.shapes.add_table(num_rows, num_cols, PH_LEFT, PH_BODY_TOP, PH_WIDTH, total_height)
    table = tbl.table

    # Set column widths if provided
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Emu(w)

    # Header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = header
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = WHITE

    # Data rows
    for row_idx, row_data in enumerate(rows):
        bg = TABLE_ALT if row_idx % 2 == 0 else None
        for col_idx, cell_val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            if bg:
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(cell_val)
            run.font.size = Pt(10)
            run.font.bold = (col_idx == 0)
            run.font.color.rgb = BLACK

    return slide


def delete_slide(prs, index):
    R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    el = prs.slides._sldIdLst[index]
    rId = el.get(f"{R_NS}id")
    if rId:
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass
    prs.slides._sldIdLst.remove(el)


def reorder_slides(prs, desired_order):
    sldIdLst = prs.slides._sldIdLst
    elements = [sldIdLst[i] for i in desired_order]
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for el in elements:
        sldIdLst.append(el)


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

prs = Presentation(BASE_PATH)

# Extract logo before any modifications
logo_path = extract_logo(prs)

# Get layouts from donor slides
content_layout = prs.slides[1].slide_layout
divider_layout  = prs.slides[2].slide_layout
table_layout    = prs.slides[3].slide_layout

# ── DONOR SLIDE 0: Cover ──────────────────────────────────────────────────
cover = prs.slides[0]
ph = cover.placeholders[0]
ph.text = "Archimedes AI\nAutomated LeanIX Population"
for shape in cover.shapes:
    if shape.has_text_frame and shape.name == "Spaker name - Dynamic":
        shape.text_frame.paragraphs[0].runs[0].text = "SAP Advisory Architecture"
    if shape.has_text_frame and shape.name == "Date - Dynamic":
        shape.text_frame.paragraphs[0].runs[0].text = "May 2026"

# ── DONOR SLIDE 1: Agenda ─────────────────────────────────────────────────
agenda_slide = prs.slides[1]
set_title(agenda_slide, "Agenda")
add_logo(agenda_slide, logo_path)
ph = agenda_slide.placeholders[10]
tf = ph.text_frame
tf.clear()
tf.word_wrap = True
agenda_items = [
    "The challenge — manual LeanIX population",
    "What is Archimedes AI",
    "How it works — 5-step pipeline",
    "Real case — Acciona",
    "Value for the Advisory Architect",
    "Getting started",
]
for i, item in enumerate(agenda_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = "\u2022  " + item
    run.font.size = Pt(15)
    run.font.color.rgb = BLACK
    p.space_before = Pt(6)
    p.space_after = Pt(6)

# ── DONOR SLIDE 2: Divider — The Challenge ───────────────────────────────
set_title(prs.slides[2], "The Challenge", is_divider=True)
add_logo(prs.slides[2], logo_path)

# ── DONOR SLIDE 3: Problem detail (table layout → use as content) ─────────
# Re-use as content slide for the problem
prob_slide = prs.slides[3]
set_title(prob_slide, "Manual LeanIX Population: A Costly Process")
add_logo(prob_slide, logo_path)
# Add bullet body manually via textbox since it's Title Only layout
from pptx.util import Inches, Pt as PPt
from pptx.util import Emu as EMU
txBox = prob_slide.shapes.add_textbox(EMU(504000), EMU(1548000), EMU(11186477), EMU(4500000))
tf = txBox.text_frame
tf.word_wrap = True
problems = [
    "2–3 days of manual work per client engagement to populate LeanIX",
    "Multiple disconnected sources: OnPrem/Cloud Excel, requirements files, PDFs, architecture diagrams",
    "Risk of inconsistency — BC coverage depends on individual knowledge of the RBA catalog",
    "Application-to-RSA mapping is manual, subjective, and error-prone",
    "Dependency on the architect who knows the SAP catalog — not scalable",
    "LeanIX stays empty or incomplete until late in the engagement",
]
for i, item in enumerate(problems):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = "\u2022  " + item
    run.font.size = Pt(14)
    run.font.color.rgb = BLACK
    p.space_before = Pt(5)
    p.space_after = Pt(5)

# ── DONOR SLIDE 4: Q&A → will be repositioned at the end ─────────────────
# We'll edit it later. First, add all new slides.

# ── NEW SLIDES ────────────────────────────────────────────────────────────

# Divider: What is Archimedes AI
s5 = add_divider_slide(prs, divider_layout, logo_path, "What is Archimedes AI")

# Slide: Overview
s6 = add_content_slide(prs, content_layout, logo_path, "Archimedes AI — Overview", [
    "AI-powered orchestrator that converts client inputs into LeanIX import files",
    "Covers the full EA cycle: AS-IS Baseline + TO-BE Target in one pipeline",
    "Engine: Claude API (claude-sonnet-4-6) + official SAP RBA and RSA catalogs",
    "RBA: 756 Business Capabilities across 22 domains",
    "RSA: 324 canonical SAP products with full hierarchy",
    "Runs locally — no infrastructure required, client data stays on your machine",
    "Output files are directly importable into LeanIX via GraphQL API",
])

# Slide: Inputs & Outputs (table)
s7 = add_table_slide(prs, table_layout, logo_path,
    "Inputs and Outputs",
    ["Input Type", "Source", "What Archimedes Extracts"],
    [
        ["OnPrem Systems Excel", "Client-provided", "Applications + hosting tags (Baseline AS-IS)"],
        ["Cloud Systems Excel", "Client-provided", "Cloud apps + solution area tags (Baseline AS-IS)"],
        ["Requirements Excel", "Client-provided", "BCs (RBA mapping) + SAP apps (RSA mapping)"],
        ["PDF documents", "Proposals, RFPs, reports", "Fact sheets: Apps, BCs, Initiatives"],
        ["Architecture diagrams", "PNG / JPG images", "Applications and IT components (Claude Vision)"],
        ["OUTPUT: baseline.xlsx", "→ LeanIX", "AS-IS applications with Baseline tags"],
        ["OUTPUT: target_leanix.xlsx", "→ LeanIX", "TO-BE: Application, BC, Initiative, ITComponent sheets"],
    ],
    col_widths=[2800000, 2800000, 5586477]
)

# Divider: How it works
s8 = add_divider_slide(prs, divider_layout, logo_path, "5-Step Pipeline")

# Slide: Steps 0–2
s9 = add_content_slide(prs, content_layout, logo_path, "Pipeline — Steps 0 to 2", [
    "Step 0 — Catalog Check: Verifies RBA/RSA catalog version and offers update",
    "Step 1 — Baseline (AS-IS): Reads OnPrem + Cloud Excel → generates baseline.xlsx with Baseline;OnPremise / Baseline;Cloud tags",
    "Step 2 — Requirements (TO-BE): Reads requirements Excel → Claude maps each row to RBA Business Capabilities and RSA products",
    "Step 2b — PDF: Extracts additional fact sheets from PDFs using Claude API (proposals, offers, reports)",
    "All steps are interactive — the pipeline prompts for each file and can skip any optional input",
])

# Slide: Steps 3–5
s10 = add_content_slide(prs, content_layout, logo_path, "Pipeline — Steps 3 to 5", [
    "Step 3 — Images: Reads architecture diagrams (.png/.jpg) → Claude Vision extracts apps and IT components",
    "Step 4 — Output generation: Produces multi-sheet LeanIX Excel (Application, BusinessCapability, Initiative, ITComponent, ReadMe)",
    "Step 5 — LeanIX import (optional): Pushes via GraphQL API in correct order: BC → ITC → App → Initiative",
    "If LEANIX_API_TOKEN is configured, import is one command: python3 run.py pipeline --client <name>",
    "Each run is idempotent — safe to re-run with updated inputs",
])

# Divider: Real Case
s11 = add_divider_slide(prs, divider_layout, logo_path, "Real Case — Acciona")

# Slide: Acciona results (table)
s12 = add_table_slide(prs, table_layout, logo_path,
    "Acciona — Pipeline Results",
    ["Phase", "Metric", "Result"],
    [
        ["Baseline AS-IS", "Total applications", "37 apps"],
        ["Baseline AS-IS", "On-Premise systems", "14 apps (Baseline;OnPremise tag)"],
        ["Baseline AS-IS", "Cloud systems", "23 apps (Baseline;Cloud tag)"],
        ["Requirements processing", "Requirements analyzed", "462 rows"],
        ["Target TO-BE", "SAP applications mapped", "7 canonical RSA apps"],
        ["Target TO-BE", "Business Capabilities (RBA)", "20 leaf BCs with parent hierarchy"],
        ["Target TO-BE", "Initiatives (by process)", "9 initiatives"],
        ["Target TO-BE", "IT Components derived", "8 ITCs (from app-to-ITC map)"],
    ],
    col_widths=[3000000, 4000000, 4186477]
)

# Slide: Acciona narrative
s13 = add_content_slide(prs, content_layout, logo_path, "Acciona — What This Means", [
    "Full LeanIX population for a real client engagement — baseline to target",
    "462 requirements processed and mapped to RBA/RSA in minutes, not days",
    "BC hierarchy automatically built from RBA: leaf nodes + parent chain",
    "Initiatives grouped by process (not by individual requirement) — ready for roadmap",
    "IT Components auto-derived from app-to-ITC static map (S/4HANA, SAC, IBP, Ariba...)",
    "Output imported into LeanIX directly — EA can start working on day 1",
])

# Divider: Value for Advisory Architect
s14 = add_divider_slide(prs, divider_layout, logo_path, "Value for Advisory Architects")

# Slide: Comparison table
s15 = add_table_slide(prs, table_layout, logo_path,
    "Without vs. With Archimedes AI",
    ["Activity", "Without Archimedes", "With Archimedes"],
    [
        ["LeanIX population", "2–3 days manual work", "Minutes (automated pipeline)"],
        ["BC coverage (RBA)", "Depends on individual knowledge", "Guaranteed — 756 BCs, 22 domains"],
        ["App-to-RSA mapping", "Manual, subjective", "Canonical resolution (324 SAP products)"],
        ["Multi-source handling", "Each source processed separately", "One orchestrator, all sources"],
        ["PDF / diagram input", "Manual extraction", "Claude API vision + NLP extraction"],
        ["LeanIX import", "Manual upload + mapping", "Direct GraphQL push, correct order"],
        ["Scalability", "Depends on who knows the catalog", "Any Advisory Architect, any client"],
    ],
    col_widths=[3500000, 3800000, 3886477]
)

# Slide: Strategic value bullets
s16 = add_content_slide(prs, content_layout, logo_path, "The Advisory Architect's Time — Refocused", [
    "From: manual data entry and catalog lookup",
    "To: strategic validation, gap analysis, and client recommendations",
    "EA value is in decisions, not in populating spreadsheets",
    "Archimedes provides a consistent, repeatable baseline from day one",
    "Enables faster time-to-value: present a populated LeanIX in the first workshop",
    "Works for any client — not tied to a specific sector or system landscape",
])

# Divider: Getting Started
s17 = add_divider_slide(prs, divider_layout, logo_path, "Getting Started")

# Slide: How to use
s18 = add_content_slide(prs, content_layout, logo_path, "Running Archimedes AI", [
    "Prerequisites: Python 3.10+, ANTHROPIC_API_KEY, optional LEANIX_API_TOKEN",
    "Clone repo: /Users/I519409/dev/archimedes-ai/",
    "Install dependencies: pip install -r requirements.txt",
    "Run the pipeline: python3 run.py pipeline --client <client_name>",
    "The pipeline is fully interactive — prompts for each input file",
    "Optional: python3 run.py enrich  (enrich only) or  run.py push  (push to LeanIX only)",
    "Output saved to: output/<client_name>/  (baseline.xlsx + target_leanix.xlsx)",
])

# ── Q&A Slide (donor slide 4) ─────────────────────────────────────────────
qa_slide = prs.slides[4]
for ph in qa_slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "Questions?"
    elif ph.placeholder_format.idx == 13:
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Let's talk about your next engagement"
        run.font.size = Pt(18)
        run.font.color.rgb = DARK_BLUE
    elif ph.placeholder_format.idx == 14:
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Archimedes AI — SAP EA Internal Tool"
        run.font.size = Pt(14)
        run.font.color.rgb = GRAY

# ── Reorder: Cover(0), Agenda(1), Div-Challenge(2), Problem(3),
#             Div-What(5=s5), Overview(6=s6), IO-table(7=s7),
#             Div-Pipeline(8=s8), Steps02(9=s9), Steps35(10=s10),
#             Div-Acciona(11=s11), AcciTable(12=s12), AcciNarr(13=s13),
#             Div-Value(14=s14), CompTable(15=s15), StrValue(16=s16),
#             Div-Start(17=s17), HowTo(18=s18), Q&A(4)
# Current indices: 0=Cover, 1=Agenda, 2=DivChallenge, 3=Problem, 4=Q&A,
#                  5..18 = new slides in order
desired_order = [0, 1, 2, 3, 5, 6, 7, 8, 9, 10, 11, 12, 13, 14, 15, 16, 17, 18, 4]
reorder_slides(prs, desired_order)

# ── Save ──────────────────────────────────────────────────────────────────
prs.save(OUTPUT_PATH)
if logo_path:
    os.unlink(logo_path)

print(f"✓ Deck saved: {OUTPUT_PATH}")
print(f"  Total slides: {len(prs.slides)}")
